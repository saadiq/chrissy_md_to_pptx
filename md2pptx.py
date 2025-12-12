#!/usr/bin/env python3
# /// script
# requires-python = ">=3.8"
# dependencies = ["python-pptx>=1.0.0"]
# ///
"""Convert markdown slides to PowerPoint presentation.

Markdown format:
    # Document Title           -> ignored (metadata)
    ## Subtitle                -> ignored (metadata)
    ---                        -> slide separator
    # SECTION N: NAME          -> section divider slide
    ## Slide N: Title          -> content slide
    | col1 | col2 |            -> table
    - item                     -> bullet point
    1. item                    -> numbered list
    **bold** *italic*          -> inline formatting
"""

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Default colors
NAVY = RGBColor(0x1C, 0x28, 0x33)
SLATE = RGBColor(0x2E, 0x40, 0x53)
SILVER = RGBColor(0xAA, 0xB7, 0xB8)
OFF_WHITE = RGBColor(0xF4, 0xF6, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def parse_markdown(md_text: str) -> list[dict]:
    """Parse markdown into slide data structures."""
    slides = []

    # Split by slide separator
    sections = re.split(r'\n---\n', md_text)

    for section in sections:
        section = section.strip()
        if not section:
            continue

        # Check for section header (# SECTION N: NAME)
        section_match = re.match(r'^#\s+SECTION\s+(\d+):\s+(.+?)(?:\s+\(\d+\s+slides?\))?$', section, re.MULTILINE)
        if section_match:
            slides.append({
                'type': 'section',
                'number': section_match.group(1),
                'title': section_match.group(2).strip()
            })
            continue

        # Check for slide content (## Slide N: Title)
        slide_match = re.match(r'^##\s+Slide\s+\d+:\s+(.+)$', section, re.MULTILINE)
        if slide_match:
            title = slide_match.group(1).strip()
            # Get content after the title line
            content_start = slide_match.end()
            content = section[content_start:].strip()

            slides.append({
                'type': 'content',
                'title': title,
                'content': content
            })
            continue

        # Check for title slide (first slide with **Title** *Subtitle* pattern)
        title_match = re.search(r'\*\*(.+?)\*\*\s*\*(.+?)\*', section)
        if title_match and 'Title Slide' in section:
            slides.append({
                'type': 'title',
                'title': title_match.group(1).strip(),
                'subtitle': title_match.group(2).strip()
            })
            continue

        # Check for document header (# Title at very beginning) - skip these
        if re.match(r'^#\s+[^#]', section) and not section_match:
            # This is document metadata, skip it
            continue

    return slides


def parse_table(text: str) -> list[list[str]]:
    """Extract table data from markdown table."""
    lines = text.strip().split('\n')
    table_data = []

    for line in lines:
        line = line.strip()
        if not line.startswith('|'):
            continue
        # Skip separator lines (|---|---|)
        if re.match(r'^\|[\s\-:]+\|', line):
            continue

        cells = [cell.strip() for cell in line.split('|')[1:-1]]
        if cells:
            table_data.append(cells)

    return table_data


def parse_bullets(text: str) -> list[dict]:
    """Extract bullet points from markdown."""
    bullets = []
    lines = text.split('\n')

    for line in lines:
        line = line.strip()
        # Checkbox bullets
        if re.match(r'^-\s*\[[ x]\]\s*', line):
            content = re.sub(r'^-\s*\[[ x]\]\s*', '', line)
            bullets.append({'text': content, 'level': 0})
        # Regular bullets
        elif re.match(r'^[-*]\s+', line):
            content = re.sub(r'^[-*]\s+', '', line)
            bullets.append({'text': content, 'level': 0})
        # Numbered items
        elif re.match(r'^\d+\.\s+', line):
            content = re.sub(r'^\d+\.\s+', '', line)
            bullets.append({'text': content, 'level': 0})
        # Indented bullets (sub-items)
        elif re.match(r'^\s+[-*]\s+', line):
            content = re.sub(r'^\s+[-*]\s+', '', line)
            bullets.append({'text': content, 'level': 1})

    return bullets


def strip_formatting(text: str) -> str:
    """Remove markdown formatting for plain text."""
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)  # bold
    text = re.sub(r'\*(.+?)\*', r'\1', text)      # italic
    text = re.sub(r'`(.+?)`', r'\1', text)        # code
    text = re.sub(r'\[(.+?)\]\((.+?)\)', r'\1', text)  # links - keep text only
    text = re.sub(r'\\(.)', r'\1', text)          # escaped chars
    text = re.sub(r'^\*{1,2}\s*', '', text)       # leading asterisks
    text = re.sub(r'\s*\*{1,2}$', '', text)       # trailing asterisks
    return text.strip()


def extract_links(text: str) -> list[dict]:
    """Extract links from markdown text."""
    return [{'text': m.group(1), 'url': m.group(2)}
            for m in re.finditer(r'\[(.+?)\]\((.+?)\)', text)]


def parse_screenshot_placeholder(text: str) -> str | None:
    """Extract screenshot placeholder description."""
    # Match: 📸 **[SCREENSHOT PLACEHOLDER]:** description
    # Handles escaped brackets \[ \] and optional markdown formatting **
    match = re.search(r'📸\s*\*{0,2}\\?\[?SCREENSHOT PLACEHOLDER\\?\]?\*{0,2}:\s*(.+)', text)
    if match:
        return strip_formatting(match.group(1))
    return None


def add_title_slide(prs, title: str, subtitle: str = ""):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Title
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = txbox.text_frame
    tf.paragraphs[0].text = strip_formatting(title)
    tf.paragraphs[0].font.size = Pt(42)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if subtitle:
        txbox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
        tf2 = txbox2.text_frame
        tf2.paragraphs[0].text = strip_formatting(subtitle)
        tf2.paragraphs[0].font.size = Pt(22)
        tf2.paragraphs[0].font.color.rgb = SILVER
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, section_num: str, title: str):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Section number
    txbox0 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.4))
    tf0 = txbox0.text_frame
    tf0.paragraphs[0].text = f"SECTION {section_num}"
    tf0.paragraphs[0].font.size = Pt(14)
    tf0.paragraphs[0].font.color.rgb = SILVER
    tf0.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(0.8))
    tf = txbox.text_frame
    tf.paragraphs[0].text = strip_formatting(title)
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title: str, content: str):
    """Add a content slide with mixed content (tables, bullets, text)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Split content into blocks and check for placeholders
    blocks = re.split(r'\n\n+', content)
    placeholders = []
    content_blocks = []

    for block in blocks:
        block = block.strip()
        if not block:
            continue
        placeholder_desc = parse_screenshot_placeholder(block)
        if placeholder_desc:
            placeholders.append(placeholder_desc)
        else:
            content_blocks.append(block)

    # Determine layout: two-column if has placeholders, full-width otherwise
    has_placeholders = len(placeholders) > 0

    # Content area dimensions
    if has_placeholders:
        content_width = Inches(5.2)  # Left side for text
        content_left = Inches(0.5)
        placeholder_left = Inches(5.9)
        placeholder_width = Inches(3.8)
    else:
        content_width = Inches(9)
        content_left = Inches(0.5)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = OFF_WHITE
    bg.line.fill.background()

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.7))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Title
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    tf = txbox.text_frame
    tf.paragraphs[0].text = strip_formatting(title)
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True

    # Add content blocks (left side if two-column)
    y_pos = Inches(0.9)

    for block in content_blocks:
        # Check if it's a table
        if '|' in block and re.search(r'^\|.+\|', block, re.MULTILINE):
            table_data = parse_table(block)
            if table_data:
                y_pos = add_table_to_slide(slide, table_data, y_pos,
                                           left=content_left, width=content_width)
                continue

        # Check if it's bullets/list
        if re.search(r'^[-*\d]', block, re.MULTILINE) or re.search(r'^\s+[-*]', block, re.MULTILINE):
            bullets = parse_bullets(block)
            if bullets:
                y_pos = add_bullets_to_slide(slide, bullets, y_pos,
                                             left=content_left, width=content_width)
                continue

        # Regular text paragraph
        if block and not block.startswith('#'):
            y_pos = add_text_to_slide(slide, block, y_pos,
                                      left=content_left, width=content_width)

    # Add placeholder(s) on right side
    if has_placeholders:
        placeholder_y = Inches(0.9)
        # Calculate height per placeholder
        available_height = Inches(4.5)
        placeholder_height = min(available_height / len(placeholders), Inches(3.5))

        for desc in placeholders:
            add_placeholder_to_slide(slide, desc, placeholder_y,
                                     left=placeholder_left,
                                     width=placeholder_width,
                                     height=placeholder_height - Inches(0.15))
            placeholder_y += placeholder_height

    return slide


def add_table_to_slide(slide, table_data: list[list[str]], y_pos,
                       left=None, width=None) -> Inches:
    """Add a table to the slide."""
    if not table_data:
        return y_pos

    left = left if left is not None else Inches(0.5)
    width = width if width is not None else Inches(9)

    rows = len(table_data)
    cols = len(table_data[0])

    # Calculate table dimensions
    table_height = Inches(0.35 * rows)
    table = slide.shapes.add_table(
        rows, cols,
        left, y_pos,
        width, table_height
    ).table

    # Style the table
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = strip_formatting(cell_text)

            # Style cell
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)

            # Header row styling
            if i == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = SLATE
            else:
                para.font.color.rgb = NAVY

    return y_pos + table_height + Inches(0.2)


def add_bullets_to_slide(slide, bullets: list[dict], y_pos,
                         left=None, width=None) -> Inches:
    """Add bullet points to the slide."""
    if not bullets:
        return y_pos

    left = left if left is not None else Inches(0.5)
    width = width if width is not None else Inches(9)

    txbox = slide.shapes.add_textbox(left, y_pos, width, Inches(3))
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = strip_formatting(bullet['text'])
        p.font.size = Pt(14)
        p.font.color.rgb = NAVY
        p.level = bullet.get('level', 0)

    # Estimate height based on bullet count
    height = Inches(0.25 * len(bullets))
    return y_pos + height + Inches(0.15)


def add_text_to_slide(slide, text: str, y_pos,
                      left=None, width=None) -> Inches:
    """Add a text paragraph to the slide."""
    left = left if left is not None else Inches(0.5)
    width = width if width is not None else Inches(9)

    txbox = slide.shapes.add_textbox(left, y_pos, width, Inches(1))
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = strip_formatting(text)
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = NAVY

    # Add hyperlinks if present
    links = extract_links(text)
    for link in links:
        # Note: python-pptx hyperlink support is limited to entire runs
        # Full implementation would require run-level formatting
        pass

    # Estimate height
    lines = len(text) // 100 + 1
    return y_pos + Inches(0.2 * lines) + Inches(0.1)


def add_placeholder_to_slide(slide, description: str, y_pos,
                             left=None, width=None, height=None):
    """Add a screenshot placeholder box to the slide."""
    left = left if left is not None else Inches(0.5)
    width = width if width is not None else Inches(9)
    height = height if height is not None else Inches(1.2)

    # Add a dashed border rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, y_pos,
        width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)  # Light gray
    shape.line.color.rgb = SILVER
    shape.line.dash_style = 2  # Dashed

    # Add placeholder text - centered in the box
    text_margin = Inches(0.15)
    txbox = slide.shapes.add_textbox(
        left + text_margin, y_pos + text_margin,
        width - (text_margin * 2), height - (text_margin * 2)
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    # Icon and label
    p = tf.paragraphs[0]
    p.text = f"[IMAGE]\n{description}"
    p.font.size = Pt(10)
    p.font.color.rgb = SLATE
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    return y_pos + height + Inches(0.15)


def convert(md_path: str, output_path: str = None, template_path: str = None):
    """Convert markdown file to PowerPoint presentation."""
    md_path = Path(md_path)

    if output_path is None:
        output_path = md_path.with_suffix('.pptx')
    else:
        output_path = Path(output_path)

    # Read markdown
    md_text = md_path.read_text(encoding='utf-8')

    # Parse slides
    slides_data = parse_markdown(md_text)

    # Create presentation
    prs = Presentation(template_path)
    if not template_path:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

    # Generate slides
    for slide_data in slides_data:
        if slide_data['type'] == 'title':
            add_title_slide(prs, slide_data['title'], slide_data.get('subtitle', ''))
        elif slide_data['type'] == 'section':
            add_section_slide(prs, slide_data['number'], slide_data['title'])
        elif slide_data['type'] == 'content':
            add_content_slide(prs, slide_data['title'], slide_data['content'])

    # Save
    prs.save(str(output_path))
    print(f"Created: {output_path}")
    print(f"Slides: {len(prs.slides)}")

    return output_path


def main():
    parser = argparse.ArgumentParser(
        description='Convert markdown slides to PowerPoint',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  uv run md2pptx.py slides.md
  uv run md2pptx.py slides.md -o presentation.pptx
  uv run md2pptx.py slides.md -t template.pptx -o output.pptx
        """
    )
    parser.add_argument('input', help='Input markdown file')
    parser.add_argument('-o', '--output', help='Output .pptx file (default: input name with .pptx)')
    parser.add_argument('-t', '--template', help='Template .pptx file for styling')

    args = parser.parse_args()
    convert(args.input, args.output, args.template)


if __name__ == '__main__':
    main()
